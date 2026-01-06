#!/usr/bin/env python3
"""
Create PowerPoint presentation from Phase 3 slide content.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import sys
from pathlib import Path

def create_presentation(output_path="Phase3_Presentation.pptx"):
    """Create PowerPoint presentation with all slides."""
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    title_color = RGBColor(31, 78, 121)  # Dark blue
    text_color = RGBColor(0, 0, 0)  # Black
    accent_color = RGBColor(0, 102, 204)  # Blue
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "On-Device Road Damage Detection\nfor Municipal Reporting"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = title_color
    title_para.alignment = PP_ALIGN.CENTER
    
    # Student name
    name_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    name_frame = name_box.text_frame
    name_frame.text = "Civan Arda Özel"
    name_para = name_frame.paragraphs[0]
    name_para.font.size = Pt(24)
    name_para.alignment = PP_ALIGN.CENTER
    
    # Supervisors
    sup_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    sup_frame = sup_box.text_frame
    sup_frame.text = "Primary Supervisor: Prof. Dr. Iftikhar Ahmed\nSecondary Supervisor: Prof. Dr. Raja Hashim Ali"
    sup_para = sup_frame.paragraphs[0]
    sup_para.font.size = Pt(18)
    sup_para.alignment = PP_ALIGN.CENTER
    
    # Dates
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.3))
    date_frame = date_box.text_frame
    date_frame.text = "December 2025 - January 2026"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(16)
    date_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem Statement
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title2 = slide2.shapes.title
    title2.text = "Problem Statement"
    title2.text_frame.paragraphs[0].font.size = Pt(36)
    title2.text_frame.paragraphs[0].font.bold = True
    title2.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "Timely road maintenance is hindered by manual inspections that miss small but safety-critical defects (e.g., cracks, potholes) under varying lighting, weather, and camera angles."
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(20)
    p2.font.bold = True
    
    # Add bullet points
    for bullet in [
        "Manual inspection limitations: miss small defects, inconsistent detection",
        "Environmental variations: lighting, weather, camera angles",
        "Safety-critical nature: small defects can lead to major hazards",
        "",
        "The Need: A lightweight, reliable object detector that can run on commodity smartphones to localize common road surface defects in real time."
    ]:
        p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.level = 0 if bullet.startswith("The Need") else 1
        if bullet.startswith("The Need"):
            p.font.bold = True
    
    # Slide 3: Research Questions
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Research Questions"
    title3.text_frame.paragraphs[0].font.size = Pt(36)
    title3.text_frame.paragraphs[0].font.bold = True
    title3.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Research Question 1 (RQ1):"
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = accent_color
    
    p3_1 = tf3.add_paragraph()
    p3_1.text = "What is the detection accuracy and performance of a YOLOv8-based model when trained on the RDD2022 dataset for road defect detection?"
    p3_1.font.size = Pt(16)
    p3_1.level = 1
    
    p3_2 = tf3.add_paragraph()
    p3_2.text = "Research Question 2 (RQ2):"
    p3_2.font.size = Pt(18)
    p3_2.font.bold = True
    p3_2.font.color.rgb = accent_color
    
    p3_3 = tf3.add_paragraph()
    p3_3.text = "Can a YOLOv8-based road defect detection model be successfully deployed on iOS mobile devices using CoreML while maintaining real-time inference performance?"
    p3_3.font.size = Pt(16)
    p3_3.level = 1
    
    p3_4 = tf3.add_paragraph()
    p3_4.text = "Research Question 3 (RQ3):"
    p3_4.font.size = Pt(18)
    p3_4.font.bold = True
    p3_4.font.color.rgb = accent_color
    
    p3_5 = tf3.add_paragraph()
    p3_5.text = "What are the usability characteristics and practical effectiveness of a mobile application for real-time road defect detection in field conditions?"
    p3_5.font.size = Pt(16)
    p3_5.level = 1
    
    # Slide 4: Literature Review Part 1
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "Literature Review / Related Work (Part 1)"
    title4.text_frame.paragraphs[0].font.size = Pt(32)
    title4.text_frame.paragraphs[0].font.bold = True
    title4.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "Object Detection Methods:"
    p4 = tf4.paragraphs[0]
    p4.font.size = Pt(18)
    p4.font.bold = True
    
    for bullet in [
        "YOLO (You Only Look Once): Real-time object detection framework",
        "YOLOv8 (2023): Latest version with improved accuracy and speed",
        "Advantages: Single-stage detection, fast inference, mobile-friendly",
        "",
        "Road Defect Detection Studies:",
        "Traditional methods: Limited accuracy (~30-40%)",
        "Large-scale CNNs: Computationally expensive, not mobile-friendly",
        "Limited research on mobile deployment",
        "",
        "Mobile ML Deployment:",
        "CoreML (Apple): Framework for on-device machine learning",
        "Edge computing trend: Moving inference to mobile devices",
        "Benefits: Privacy, low latency, offline capability"
    ]:
        p = tf4.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        if bullet.endswith(":") or bullet == "":
            p.level = 0
            if bullet.endswith(":"):
                p.font.bold = True
        else:
            p.level = 1
    
    # Slide 5: Literature Review Part 2
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "Literature Review / Related Work (Part 2)"
    title5.text_frame.paragraphs[0].font.size = Pt(32)
    title5.text_frame.paragraphs[0].font.bold = True
    title5.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "RDD2022 Dataset:"
    p5 = tf5.paragraphs[0]
    p5.font.size = Pt(18)
    p5.font.bold = True
    
    for bullet in [
        "19,089 training images, 3,579 validation images",
        "6 defect classes: D00, D01, D10, D11, D20, D40",
        "Standard benchmark for road defect detection research",
        "",
        "Performance Benchmarks:",
        "Baseline models: 30-45% mAP@0.5",
        "State-of-the-art: 60-70% mAP@0.5 (larger models)",
        "Gap: Need for mobile-deployable models with >45% mAP@0.5"
    ]:
        p = tf5.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        if bullet.endswith(":") or bullet == "":
            p.level = 0
            if bullet.endswith(":"):
                p.font.bold = True
        else:
            p.level = 1
    
    # Slide 6: Methodology
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "Methodology"
    title6.text_frame.paragraphs[0].font.size = Pt(36)
    title6.text_frame.paragraphs[0].font.bold = True
    title6.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "Dataset:"
    p6 = tf6.paragraphs[0]
    p6.font.size = Pt(18)
    p6.font.bold = True
    
    for bullet in [
        "RDD2022: 19,089 training images, 3,579 validation images",
        "6 defect classes, YOLO format conversion",
        "",
        "Model Architecture:",
        "YOLOv8s (Small variant): 11M parameters",
        "Pre-trained on COCO, fine-tuned on RDD2022",
        "Image size: 640x640 pixels",
        "",
        "Training Strategy:",
        "200 epochs with cosine learning rate schedule",
        "Data augmentation: mosaic, mixup, HSV, rotation, scaling",
        "Batch size: 4, Device: MPS (macOS)",
        "",
        "Mobile Deployment:",
        "CoreML conversion for iOS",
        "SwiftUI native application",
        "Real-time camera inference",
        "GPS location tagging",
        "Municipal reporting system (25+ countries)"
    ]:
        p = tf6.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        if bullet.endswith(":") or bullet == "":
            p.level = 0
            if bullet.endswith(":"):
                p.font.bold = True
        else:
            p.level = 1
    
    # Slide 7: Results RQ1
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    title7.text = "Results - Research Question 1"
    title7.text_frame.paragraphs[0].font.size = Pt(32)
    title7.text_frame.paragraphs[0].font.bold = True
    title7.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.text = "RQ1: Detection Accuracy and Performance"
    p7 = tf7.paragraphs[0]
    p7.font.size = Pt(20)
    p7.font.bold = True
    p7.font.color.rgb = accent_color
    
    for bullet in [
        "",
        "Best Performance Achieved:",
        "• mAP@0.5: 49.19% (Epoch 119)",
        "• mAP@0.5:0.95: 20.58% (Epoch 87)",
        "• Precision: 60.82% (Epoch 120)",
        "• Recall: 47.22% (Epoch 152)",
        "",
        "Progress:",
        "• Started: 16.01% mAP@0.5 (Epoch 1)",
        "• Best: 49.19% mAP@0.5 (Epoch 119)",
        "• Improvement: 207% increase",
        "• Status: 179/200 epochs completed (90%)",
        "",
        "Comparison:",
        "✅ Exceeds baseline models (30-45%)",
        "✅ Comparable to YOLOv8s expectations (45-55%)",
        "✅ Strong mobile-optimized performance",
        "",
        "Answer to RQ1:",
        "✅ Yes - YOLOv8s achieves effective road defect detection with 49.19% mAP@0.5 and 60.82% precision, exceeding baseline methods."
    ]:
        p = tf7.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        if bullet.startswith("•") or bullet.startswith("✅"):
            p.level = 1
        elif bullet.endswith(":") or bullet == "":
            p.level = 0
            if bullet.endswith(":"):
                p.font.bold = True
        else:
            p.level = 0
            if "Answer" in bullet:
                p.font.bold = True
                p.font.color.rgb = accent_color
    
    # Slide 8: Results RQ2
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    title8.text = "Results - Research Question 2"
    title8.text_frame.paragraphs[0].font.size = Pt(32)
    title8.text_frame.paragraphs[0].font.bold = True
    title8.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content8 = slide8.placeholders[1]
    tf8 = content8.text_frame
    tf8.text = "RQ2: Mobile Deployment Feasibility"
    p8 = tf8.paragraphs[0]
    p8.font.size = Pt(20)
    p8.font.bold = True
    p8.font.color.rgb = accent_color
    
    for bullet in [
        "",
        "CoreML Conversion:",
        "✅ Conversion successful",
        "✅ iOS Neural Engine compatible",
        "✅ Optimized for on-device inference",
        "",
        "Model Specifications:",
        "• Model Size: ~43 MB (acceptable for mobile)",
        "• Inference Time: <50ms on iPhone 12+",
        "• Format: CoreML (.mlmodel), iOS 14.0+",
        "",
        "iOS App Implementation:",
        "✅ Native SwiftUI application",
        "✅ Real-time camera inference",
        "✅ CoreML integration complete",
        "✅ GPS location tagging",
        "✅ Detection history management",
        "✅ Municipal reporting system (25+ countries)",
        "  - Phone numbers and email addresses",
        "  - One-tap calling and email",
        "  - Automatic authority identification",
        "",
        "Answer to RQ2:",
        "✅ Yes - Successfully deployed to iOS using CoreML with real-time inference <50ms and complete application features."
    ]:
        p = tf8.add_paragraph()
        p.text = bullet
        p.font.size = Pt(13)
        if bullet.startswith("  -"):
            p.level = 2
        elif bullet.startswith("•") or bullet.startswith("✅"):
            p.level = 1
        elif bullet.endswith(":") or bullet == "":
            p.level = 0
            if bullet.endswith(":"):
                p.font.bold = True
        else:
            p.level = 0
            if "Answer" in bullet:
                p.font.bold = True
                p.font.color.rgb = accent_color
    
    # Slide 9: Results RQ3
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    title9.text = "Results - Research Question 3"
    title9.text_frame.paragraphs[0].font.size = Pt(32)
    title9.text_frame.paragraphs[0].font.bold = True
    title9.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content9 = slide9.placeholders[1]
    tf9 = content9.text_frame
    tf9.text = "RQ3: Usability and Practical Effectiveness"
    p9 = tf9.paragraphs[0]
    p9.font.size = Pt(20)
    p9.font.bold = True
    p9.font.color.rgb = accent_color
    
    for bullet in [
        "",
        "iOS App Features:",
        "✅ Real-time defect detection via camera",
        "✅ Photo library import",
        "✅ Bounding box visualization with confidence scores",
        "✅ Per-class color coding (6 defect types)",
        "✅ GPS location tagging",
        "✅ Detection history management",
        "✅ Municipal reporting (25+ countries)",
        "  - Contact database: phone numbers and emails",
        "  - Automatic authority identification",
        "  - One-tap calling and email",
        "",
        "User Experience:",
        "• Modern SwiftUI interface",
        "• Intuitive navigation",
        "• Professional UI/UX design",
        "• Responsive interactions",
        "",
        "Practical Utility:",
        "• Field-ready: Works on-site with smartphones",
        "• Offline capability: No internet required",
        "• Direct reporting: Contact authorities from app",
        "• Historical tracking: Maintains detection records",
        "",
        "Answer to RQ3:",
        "✅ Yes - Application demonstrates strong usability and practical effectiveness with intuitive interface, real-time detection, and comprehensive reporting capabilities."
    ]:
        p = tf9.add_paragraph()
        p.text = bullet
        p.font.size = Pt(13)
        if bullet.startswith("  -"):
            p.level = 2
        elif bullet.startswith("•") or bullet.startswith("✅"):
            p.level = 1
        elif bullet.endswith(":") or bullet == "":
            p.level = 0
            if bullet.endswith(":"):
                p.font.bold = True
        else:
            p.level = 0
            if "Answer" in bullet:
                p.font.bold = True
                p.font.color.rgb = accent_color
    
    # Slide 10: Conclusion
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    title10.text = "Conclusion"
    title10.text_frame.paragraphs[0].font.size = Pt(36)
    title10.text_frame.paragraphs[0].font.bold = True
    title10.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content10 = slide10.placeholders[1]
    tf10 = content10.text_frame
    tf10.text = "Answers to Research Questions:"
    p10 = tf10.paragraphs[0]
    p10.font.size = Pt(18)
    p10.font.bold = True
    
    for bullet in [
        "",
        "RQ1: ✅ Yes - 49.19% mAP@0.5, 60.82% precision, exceeds baseline methods",
        "",
        "RQ2: ✅ Yes - Successfully deployed to iOS, <50ms inference, complete features",
        "",
        "RQ3: ✅ Yes - Strong usability, practical effectiveness, ready for field deployment",
        "",
        "Key Contributions:",
        "1. Trained YOLOv8s model: 49.19% mAP@0.5, 60.82% precision",
        "2. Successfully deployed to iOS: CoreML conversion, real-time inference",
        "3. Complete iOS application: Real-time detection, GPS tagging, municipal reporting",
        "",
        "Key Achievements:",
        "✅ End-to-end solution: Dataset → Training → Mobile Deployment",
        "✅ Above-baseline performance (49.19% vs 30-45% baseline)",
        "✅ Production-ready iOS application",
        "✅ Comprehensive documentation",
        "",
        "Impact:",
        "• Practical solution for road maintenance",
        "• Demonstrates mobile ML deployment feasibility",
        "• Contributes to infrastructure monitoring technology",
        "• Enables field workers to perform on-site detection"
    ]:
        p = tf10.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        if bullet.startswith("•") or bullet.startswith("✅") or bullet.startswith(("1.", "2.", "3.")):
            p.level = 1
        elif bullet.endswith(":") or bullet == "":
            p.level = 0
            if bullet.endswith(":"):
                p.font.bold = True
        else:
            p.level = 0
    
    # Save presentation
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    output_file = sys.argv[1] if len(sys.argv) > 1 else "Phase3_Presentation.pptx"
    try:
        create_presentation(output_file)
    except ImportError:
        print("❌ Error: python-pptx library not installed.")
        print("   Please install it with: pip install python-pptx")
        sys.exit(1)
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        sys.exit(1)

